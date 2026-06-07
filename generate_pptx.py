#!/usr/bin/env python3
"""
NeuroBalance — 4-slide presentation generator
Real data from FINAL_PIPELINE.ipynb
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy, io

NB_DIR = '/Users/amangeldimadina/IdeaProjects/NeuroBalanceBackend/ml/notebooks/'

# ── Palette ─────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x0D, 0x1B, 0x2A)   # почти чёрный синий
ACCENT1    = RGBColor(0x00, 0xC8, 0xFF)   # яркий голубой
ACCENT2    = RGBColor(0x7C, 0x4D, 0xFF)   # фиолетовый
ACCENT3    = RGBColor(0x00, 0xE6, 0x96)   # мятный
ACCENT4    = RGBColor(0xFF, 0x6B, 0x6B)   # красный/оранжевый
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xBE, 0xC5)
CARD_BG    = RGBColor(0x16, 0x2B, 0x3F)   # тёмная карточка

W = Inches(13.33)   # 16:9 widescreen
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank = prs.slide_layouts[6]   # полностью пустой layout

# ════════════════════════════════════════════════════════════════
# Helpers
# ════════════════════════════════════════════════════════════════

def rgb(r, g, b):
    return RGBColor(r, g, b)

def add_rect(slide, x, y, w, h, fill=None, alpha=None,
             line_color=None, line_width=Pt(1), radius=0):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE=1
    fill_obj = shape.fill
    if fill:
        fill_obj.solid()
        fill_obj.fore_color.rgb = fill
    else:
        fill_obj.background()
    line = shape.line
    if line_color:
        line.color.rgb = line_color
        line.width = line_width
    else:
        line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tf_box = slide.shapes.add_textbox(x, y, w, h)
    tf = tf_box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tf_box

def add_image(slide, path, x, y, w, h=None):
    try:
        if h:
            pic = slide.shapes.add_picture(path, x, y, w, h)
        else:
            pic = slide.shapes.add_picture(path, x, y, w)
        return pic
    except Exception as e:
        print(f"Image error {path}: {e}")
        return None

def dark_bg(slide):
    add_rect(slide, 0, 0, W, H, fill=DARK_BG)

def slide_header(slide, title, subtitle=None, accent=ACCENT1):
    """Горизонтальная полоса заголовка сверху"""
    bar = add_rect(slide, 0, 0, W, Inches(1.1), fill=CARD_BG)
    # Акцентная линия снизу заголовка
    add_rect(slide, 0, Inches(1.1), W, Pt(3), fill=accent)
    add_text(slide, title,
             Inches(0.4), Inches(0.12), Inches(9), Inches(0.7),
             size=32, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.4), Inches(0.68), Inches(10), Inches(0.38),
                 size=14, color=LIGHT_GRAY)

def badge(slide, text, x, y, fill=ACCENT1, text_color=DARK_BG, size=13):
    w = Inches(1.5); h = Inches(0.35)
    add_rect(slide, x, y, w, h, fill=fill)
    add_text(slide, text, x, y, w, h,
             size=size, bold=True, color=text_color, align=PP_ALIGN.CENTER)

def metric_card(slide, label, value, unit, x, y, accent=ACCENT1):
    cw = Inches(1.9); ch = Inches(1.1)
    add_rect(slide, x, y, cw, ch, fill=CARD_BG,
             line_color=accent, line_width=Pt(1.5))
    add_text(slide, value, x, y+Pt(6), cw, Inches(0.55),
             size=28, bold=True, color=accent, align=PP_ALIGN.CENTER)
    add_text(slide, unit,
             x, y+Inches(0.52), cw, Inches(0.25),
             size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, label,
             x, y+Inches(0.75), cw, Inches(0.3),
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 1 — ML Модели: Когнитивный скор + Классификатор состояния
# ════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank)
dark_bg(s1)
slide_header(s1,
    "ML Модели: Когнитивный скор и Состояние",
    "XGBoost Regressor + XGBoost Classifier  ·  Датасет: 80 000 записей  ·  Split 70/10/20",
    accent=ACCENT1)

# ── Левая половина: Когнитивный скор ────────────────────────────
add_rect(s1, Inches(0.3), Inches(1.25), Inches(6.0), Inches(5.95),
         fill=CARD_BG, line_color=ACCENT1, line_width=Pt(1.5))

add_text(s1, "① XGBoost Regressor — Когнитивный Скор",
         Inches(0.45), Inches(1.32), Inches(5.7), Inches(0.4),
         size=15, bold=True, color=ACCENT1)

add_text(s1,
    "Задача: предсказать непрерывное значение\n"
    "когнитивного скора от 0 до 100 на основе\n"
    "17 признаков образа жизни.",
    Inches(0.45), Inches(1.72), Inches(5.7), Inches(0.7),
    size=11, color=LIGHT_GRAY)

# Метрики регрессии
metric_card(s1, "Test R²",    "0.9997", "из 1.0",  Inches(0.45), Inches(2.45), ACCENT1)
metric_card(s1, "Test RMSE",  "0.370",  "баллов",  Inches(2.45), Inches(2.45), ACCENT3)
metric_card(s1, "Val RMSE",   "0.360",  "баллов",  Inches(4.45), Inches(2.45), ACCENT2)

add_text(s1,
    "R² = 0.9997 означает, что модель объясняет\n"
    "99.97% дисперсии когнитивного скора.\n"
    "RMSE = 0.370 — ошибка менее 0.4 балла из 100.",
    Inches(0.45), Inches(3.62), Inches(5.7), Inches(0.65),
    size=10.5, color=LIGHT_GRAY, italic=True)

# Сравнение с Random Forest
add_rect(s1, Inches(0.45), Inches(4.32), Inches(5.7), Inches(0.72),
         fill=rgb(0x1a,0x2e,0x40), line_color=ACCENT3, line_width=Pt(1))
add_text(s1, "Сравнение: XGBoost vs Random Forest",
         Inches(0.55), Inches(4.34), Inches(5.5), Inches(0.28),
         size=11, bold=True, color=ACCENT3)

# Полоски сравнения
def bar_compare(slide, x, y, label, val1, val2, lbl1, lbl2, w_total=Inches(4.5)):
    add_text(slide, label, x, y, Inches(1.5), Inches(0.22), size=10, color=WHITE)
    # XGBoost bar
    bw1 = w_total * val1
    add_rect(slide, x + Inches(1.55), y + Pt(2), bw1, Inches(0.18), fill=ACCENT1)
    add_text(slide, lbl1, x + Inches(1.55) + bw1 + Pt(4), y, Inches(1.0), Inches(0.22),
             size=9, bold=True, color=ACCENT1)
    # RF bar
    bw2 = w_total * val2
    add_rect(slide, x + Inches(1.55), y + Inches(0.22), bw2, Inches(0.18), fill=ACCENT4)
    add_text(slide, lbl2, x + Inches(1.55) + bw2 + Pt(4), y + Inches(0.22), Inches(1.0), Inches(0.22),
             size=9, bold=True, color=ACCENT4)

bar_compare(s1, Inches(0.55), Inches(4.62),
            "Test R²:",
            0.9997, 0.9983,
            "XGB: 0.9997", "RF:  0.9983",
            w_total=Inches(3.5))

add_text(s1,
    "17 признаков: 10 исходных + 7 engineered\n"
    "(CFI, Sleep_Debt, Memory_Efficiency,\n"
    "Lifestyle_Balance, Sleep×Exercise, …)",
    Inches(0.45), Inches(5.1), Inches(5.7), Inches(0.65),
    size=10, color=LIGHT_GRAY)

add_text(s1, "Поч. XGBoost: L1+L2 регуляризация,\nearly stopping, 0.37 RMSE vs 0.93 у RF",
         Inches(0.45), Inches(5.8), Inches(5.7), Inches(0.35),
         size=9.5, color=ACCENT3, italic=True)

# ── Правая половина: Классификатор ──────────────────────────────
add_rect(s1, Inches(6.65), Inches(1.25), Inches(6.35), Inches(5.95),
         fill=CARD_BG, line_color=ACCENT2, line_width=Pt(1.5))

add_text(s1, "② XGBoost Classifier — Когнитивное Состояние",
         Inches(6.8), Inches(1.32), Inches(6.0), Inches(0.4),
         size=15, bold=True, color=ACCENT2)

add_text(s1,
    "Задача: 3-классовая классификация состояния.\n"
    "HIGH (>75 баллов) · MEDIUM (50–75) · LOW (<50)",
    Inches(6.8), Inches(1.72), Inches(6.0), Inches(0.55),
    size=11, color=LIGHT_GRAY)

# Метрики классификации
metric_card(s1, "Val Accuracy",  "97.94%", "на val set",  Inches(6.8),  Inches(2.32), ACCENT2)
metric_card(s1, "Test Accuracy", "99.99%", "на test set", Inches(8.85), Inches(2.32), ACCENT3)
metric_card(s1, "Train size",    "56 000", "записей",     Inches(10.9), Inches(2.32), ACCENT1)

# Classification Report table
add_text(s1, "Classification Report (Val Set — 8 000 образцов):",
         Inches(6.8), Inches(3.5), Inches(6.0), Inches(0.3),
         size=11, bold=True, color=WHITE)

headers = ["Класс",    "Precision", "Recall", "F1",   "Support"]
rows = [
    ["HIGH",   "0.99",    "0.98",   "0.98",  "2 637"],
    ["LOW",    "0.98",    "0.98",   "0.98",  "1 903"],
    ["MEDIUM", "0.97",    "0.98",   "0.98",  "3 460"],
]
col_x = [Inches(6.8), Inches(8.1), Inches(9.15), Inches(10.1), Inches(11.0)]
col_w = [Inches(1.2), Inches(0.95), Inches(0.85), Inches(0.8), Inches(1.0)]
row_h = Inches(0.32)
header_y = Inches(3.82)

# Header row
for ci, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
    add_rect(s1, cx, header_y, cw, row_h, fill=ACCENT2)
    add_text(s1, hdr, cx+Pt(4), header_y+Pt(4), cw, row_h,
             size=10, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)

class_colors = [ACCENT3, ACCENT1, ACCENT4]
for ri, row in enumerate(rows):
    ry = header_y + row_h * (ri+1)
    row_bg = rgb(0x16,0x2b,0x40) if ri%2==0 else rgb(0x1a,0x30,0x47)
    for ci, (val, cx, cw) in enumerate(zip(row, col_x, col_w)):
        add_rect(s1, cx, ry, cw, row_h, fill=row_bg,
                 line_color=rgb(0x2a,0x3f,0x55), line_width=Pt(0.5))
        col_clr = class_colors[ri] if ci==0 else WHITE
        add_text(s1, val, cx+Pt(4), ry+Pt(4), cw, row_h,
                 size=10, bold=(ci==0), color=col_clr, align=PP_ALIGN.CENTER)

# Распределение классов — мини бар-чарт
add_text(s1, "Распределение классов в датасете:",
         Inches(6.8), Inches(5.05), Inches(6.0), Inches(0.28),
         size=10.5, bold=True, color=WHITE)

dist = [("MEDIUM", 34595, ACCENT4), ("HIGH", 26374, ACCENT3), ("LOW", 19031, ACCENT1)]
total = 80000
for bi, (lbl, cnt, col) in enumerate(dist):
    bx = Inches(6.8)
    by = Inches(5.35) + Inches(0.38) * bi
    bw_full = Inches(5.5)
    bw = bw_full * (cnt / total)
    add_rect(s1, bx, by, bw, Inches(0.28), fill=col)
    pct = f"{cnt/total*100:.1f}%  {cnt:,}"
    add_text(s1, f"{lbl}: {pct}", bx + bw + Pt(6), by, Inches(2.5), Inches(0.28),
             size=10, color=col)

add_text(s1,
    "Модель использует когнитивный скор для адаптивной\n"
    "сложности игр: >75 → HARD, 50-75 → MEDIUM, <50 → EASY",
    Inches(6.8), Inches(6.55), Inches(6.0), Inches(0.55),
    size=9.5, color=LIGHT_GRAY, italic=True)

# Slide number
add_text(s1, "1 / 4", Inches(12.5), Inches(7.2), Inches(0.7), Inches(0.25),
         size=10, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════════
# SLIDE 2 — CatBoost + Health Metrics (M-Rest / M-Ready / M-Balance)
# ════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank)
dark_bg(s2)
slide_header(s2,
    "ML Модели: Рекомендации и Метрики Здоровья",
    "CatBoost Regressor (рекомендации)  ·  XGBoost ×3 (M-Rest / M-Ready / M-Balance)  ·  12 000 записей",
    accent=ACCENT2)

# ── Левая: CatBoost ─────────────────────────────────────────────
add_rect(s2, Inches(0.3), Inches(1.25), Inches(6.0), Inches(5.95),
         fill=CARD_BG, line_color=ACCENT2, line_width=Pt(1.5))

add_text(s2, "③ CatBoost Regressor — Персональные Рекомендации",
         Inches(0.45), Inches(1.32), Inches(5.7), Inches(0.4),
         size=14, bold=True, color=ACCENT2)

add_text(s2,
    "Задача: для каждого из 4 типов рекомендаций предсказать\n"
    "predicted_improvement — прирост когнитивного скора.\n"
    "Выбирает ТОП-3 с макс. improvement без дубликатов.",
    Inches(0.45), Inches(1.72), Inches(5.7), Inches(0.65),
    size=10.5, color=LIGHT_GRAY)

metric_card(s2, "Test R²",   "0.9656", "из 1.0",      Inches(0.45), Inches(2.42), ACCENT2)
metric_card(s2, "RMSE",      "1.806",  "баллов",       Inches(2.45), Inches(2.42), ACCENT3)
metric_card(s2, "MAE",       "1.20",   "баллов",       Inches(4.45), Inches(2.42), ACCENT1)

add_text(s2,
    "RMSE = 1.806 → ошибка предсказания\n"
    "улучшения <2 баллов из 100. Обучено на ~50 000 сценариев\n"
    "(4 типа × разные величины change × 12 500 профилей).",
    Inches(0.45), Inches(3.6), Inches(5.7), Inches(0.65),
    size=10, color=LIGHT_GRAY, italic=True)

# 4 типа рекомендаций
add_text(s2, "4 типа рекомендаций (выходы модели):",
         Inches(0.45), Inches(4.3), Inches(5.7), Inches(0.3),
         size=11, bold=True, color=WHITE)

rec_items = [
    ("🌙", "SLEEP_INCREASE",   "Оптимизировать сон",       ACCENT1),
    ("🧘", "STRESS_DECREASE",  "Снизить стресс",            ACCENT2),
    ("💪", "EXERCISE_INCREASE","Увеличить активность",      ACCENT3),
    ("📱", "SCREEN_DECREASE",  "Сократить экранное время",  ACCENT4),
]
for ri, (emoji, rtype, rdesc, rcol) in enumerate(rec_items):
    ry = Inches(4.65) + Inches(0.4) * ri
    add_rect(s2, Inches(0.45), ry, Inches(5.7), Inches(0.33),
             fill=rgb(0x1a,0x2e,0x40), line_color=rcol, line_width=Pt(1))
    add_text(s2, f"{emoji} {rtype}",
             Inches(0.55), ry+Pt(3), Inches(2.5), Inches(0.28),
             size=10, bold=True, color=rcol)
    add_text(s2, f"— {rdesc}",
             Inches(3.1), ry+Pt(3), Inches(2.9), Inches(0.28),
             size=10, color=LIGHT_GRAY)

add_text(s2,
    "Почему CatBoost: нативные категориальные признаки\n"
    "(action_type), Ordered Target Statistics — нет target leakage.\n"
    "Вердикт ноутбука: «Модель обучена ВЕЛИКОЛЕПНО!»",
    Inches(0.45), Inches(6.35), Inches(5.7), Inches(0.7),
    size=9.5, color=ACCENT3, italic=True)

# ── Правая: Health Metrics ───────────────────────────────────────
add_rect(s2, Inches(6.65), Inches(1.25), Inches(6.35), Inches(5.95),
         fill=CARD_BG, line_color=ACCENT3, line_width=Pt(1.5))

add_text(s2, "④⑤⑥ XGBoost ×3 — M-Rest / M-Ready / M-Balance",
         Inches(6.8), Inches(1.32), Inches(6.0), Inches(0.4),
         size=14, bold=True, color=ACCENT3)

add_text(s2,
    "Три отдельные модели, каждая специализирована\n"
    "под свою метрику. Обучены на 12 000 синтетических\n"
    "записей с физиологически валидированными формулами.",
    Inches(6.8), Inches(1.72), Inches(6.0), Inches(0.6),
    size=10.5, color=LIGHT_GRAY)

# Таблица метрик
add_text(s2, "Результаты обучения (test set):",
         Inches(6.8), Inches(2.36), Inches(6.0), Inches(0.3),
         size=11, bold=True, color=WHITE)

hm_headers = ["Метрика", "R²", "MAE", "RMSE", "Ср. значение"]
hm_rows = [
    ["M-Rest",    "0.9170", "2.01", "2.52", "85.2 ± 8.7"],
    ["M-Ready",   "0.9713", "2.02", "2.56", "65.6 ± 14.8"],
    ["M-Balance", "0.9674", "2.13", "2.66", "60.0 ± 14.7"],
]
hm_col_x = [Inches(6.8), Inches(8.25), Inches(9.3), Inches(10.25), Inches(11.2)]
hm_col_w = [Inches(1.35), Inches(0.95), Inches(0.85), Inches(0.85), Inches(1.55)]
hm_colors = [ACCENT1, ACCENT2, ACCENT3]
hdr_y2 = Inches(2.7)

for ci, (hdr, cx, cw) in enumerate(zip(hm_headers, hm_col_x, hm_col_w)):
    add_rect(s2, cx, hdr_y2, cw, row_h, fill=ACCENT3)
    add_text(s2, hdr, cx+Pt(3), hdr_y2+Pt(4), cw, row_h,
             size=10, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)

for ri, (row, clr) in enumerate(zip(hm_rows, hm_colors)):
    ry = hdr_y2 + row_h * (ri+1)
    for ci, (val, cx, cw) in enumerate(zip(row, hm_col_x, hm_col_w)):
        bg = rgb(0x16,0x2b,0x40) if ri%2==0 else rgb(0x1a,0x30,0x47)
        add_rect(s2, cx, ry, cw, row_h, fill=bg,
                 line_color=rgb(0x2a,0x3f,0x55), line_width=Pt(0.5))
        c = clr if ci==0 else WHITE
        add_text(s2, val, cx+Pt(3), ry+Pt(4), cw, row_h,
                 size=10.5, bold=(ci==0), color=c, align=PP_ALIGN.CENTER)

# Формулы
add_text(s2, "Формулы (детерминированный fallback):",
         Inches(6.8), Inches(4.05), Inches(6.0), Inches(0.3),
         size=11, bold=True, color=WHITE)

formulas = [
    ("M-Rest",    ACCENT1,  "min(hours/8,1)×50 + quality×30 + rested×20"),
    ("M-Ready",   ACCENT2,  "(energy/10)×40 + (mood/5)×30 + M-Rest×0.30"),
    ("M-Balance", ACCENT3,  "((10-stress)/9)×40 + habits×40 + mood×20"),
]
for fi, (name, clr, formula) in enumerate(formulas):
    fy = Inches(4.38) + Inches(0.5) * fi
    add_rect(s2, Inches(6.8), fy, Inches(6.0), Inches(0.43),
             fill=rgb(0x1a,0x2e,0x40), line_color=clr, line_width=Pt(1.5))
    add_text(s2, name, Inches(6.9), fy+Pt(5), Inches(1.05), Inches(0.3),
             size=10.5, bold=True, color=clr)
    add_text(s2, formula, Inches(8.0), fy+Pt(5), Inches(4.7), Inches(0.3),
             size=9.5, color=WHITE)

# Visual progress bars for R²
add_text(s2, "R² визуально (чем длиннее — тем точнее):",
         Inches(6.8), Inches(5.96), Inches(6.0), Inches(0.28),
         size=10.5, bold=True, color=WHITE)

r2_data = [("M-Rest", 0.9170, ACCENT1), ("M-Ready", 0.9713, ACCENT2), ("M-Balance", 0.9674, ACCENT3)]
for ri, (name, r2val, clr) in enumerate(r2_data):
    by = Inches(6.27) + Inches(0.33)*ri
    max_w = Inches(3.8)
    add_rect(s2, Inches(6.8), by, max_w, Inches(0.22),
             fill=rgb(0x25,0x3a,0x50))
    add_rect(s2, Inches(6.8), by, max_w*r2val, Inches(0.22), fill=clr)
    add_text(s2, f"{name}: {r2val:.4f}",
             Inches(10.7), by, Inches(2.2), Inches(0.22),
             size=10, bold=True, color=clr)

add_text(s2, "2 / 4", Inches(12.5), Inches(7.2), Inches(0.7), Inches(0.25),
         size=10, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════════
# SLIDE 3 — Системная Архитектура
# ════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank)
dark_bg(s3)
slide_header(s3,
    "Системная Архитектура NeuroBalance",
    "Microservices · Docker Compose 10 containers · Database-per-Service · Async Kafka Events",
    accent=ACCENT3)

# ── iOS Client ──────────────────────────────────────────────────
ios_x = Inches(0.25); ios_y = Inches(1.35)
ios_w = Inches(2.3);  ios_h = Inches(1.3)
add_rect(s3, ios_x, ios_y, ios_w, ios_h,
         fill=rgb(0x16,0x2b,0x40), line_color=ACCENT1, line_width=Pt(2))
add_text(s3, "📱 iOS App", ios_x, ios_y+Pt(8), ios_w, Inches(0.35),
         size=14, bold=True, color=ACCENT1, align=PP_ALIGN.CENTER)
add_text(s3, "SwiftUI · MVVM\nURLSession async/await\nKeychain JWT storage",
         ios_x+Pt(6), ios_y+Inches(0.42), ios_w-Pt(12), Inches(0.8),
         size=9, color=LIGHT_GRAY)

# ── Arrow iOS → Services ──────────────────────────────────────
add_text(s3, "HTTPS\nJWT Bearer",
         Inches(2.62), Inches(1.72), Inches(0.75), Inches(0.5),
         size=8, color=ACCENT1, align=PP_ALIGN.CENTER)
add_rect(s3, Inches(2.6), Inches(1.98), Inches(0.65), Pt(2), fill=ACCENT1)

# ── 3 Java Services ─────────────────────────────────────────────
svc_configs = [
    ("NBAuthService\n:8081", "JWT · BCrypt\nRegistration\nOnboarding\nProfile", ACCENT1,
     Inches(3.35), Inches(1.35)),
    ("NBCheckinService\n:8082", "Checkins · Mood\nSleep · Tasks\nGames · XP\nML Results", ACCENT2,
     Inches(5.85), Inches(1.35)),
    ("NoteAI-backend\n:8083", "Journal · Notes\nGroq Llama 3.3\n70b AI Analysis\nAuto-save PATCH", ACCENT4,
     Inches(8.35), Inches(1.35)),
]
svc_w = Inches(2.3); svc_h = Inches(1.85)
for (title, desc, clr, sx, sy) in svc_configs:
    add_rect(s3, sx, sy, svc_w, svc_h,
             fill=CARD_BG, line_color=clr, line_width=Pt(2))
    add_text(s3, title, sx+Pt(5), sy+Pt(6), svc_w-Pt(10), Inches(0.45),
             size=12, bold=True, color=clr)
    add_text(s3, desc, sx+Pt(6), sy+Inches(0.5), svc_w-Pt(12), Inches(1.25),
             size=9, color=LIGHT_GRAY)

# Spring Boot badge на каждом
for (_, _, clr, sx, sy) in svc_configs:
    badge(s3, "Spring Boot 4", sx + Inches(0.1), sy + svc_h - Inches(0.02),
          fill=clr, text_color=DARK_BG, size=8)

# ── Arrow Services → Kafka ──────────────────────────────────────
kafka_x = Inches(5.85); kafka_y = Inches(3.6)
kafka_w = Inches(4.75); kafka_h = Inches(1.0)

# vertical arrows down
for sx_arrow in [Inches(4.5), Inches(6.65), Inches(9.15)]:
    for dot_y in [Inches(3.22), Inches(3.38), Inches(3.54)]:
        add_rect(s3, sx_arrow, dot_y, Pt(3), Pt(3), fill=ACCENT2)

add_rect(s3, kafka_x, kafka_y, kafka_w, kafka_h,
         fill=rgb(0x1a,0x2e,0x40), line_color=ACCENT3, line_width=Pt(2))
add_text(s3, "⚡ Apache Kafka  :29092",
         kafka_x+Pt(8), kafka_y+Pt(5), kafka_w-Pt(16), Inches(0.35),
         size=13, bold=True, color=ACCENT3)
add_text(s3, "checkin.created  ·  sleep.logged  ·  game.completed  ·  character.leveled-up",
         kafka_x+Pt(8), kafka_y+Inches(0.42), kafka_w-Pt(16), Inches(0.28),
         size=8.5, color=LIGHT_GRAY)
add_text(s3, "Partition key = userId → гарантия порядка событий | Replication factor = 1 | 3 partitions/topic",
         kafka_x+Pt(8), kafka_y+Inches(0.67), kafka_w-Pt(16), Inches(0.28),
         size=7.5, color=rgb(0x70,0x80,0x90), italic=True)

# ── ML Service ──────────────────────────────────────────────────
ml_x = Inches(5.85); ml_y = Inches(4.85)
ml_w = Inches(2.2);  ml_h = Inches(1.7)
add_rect(s3, ml_x, ml_y, ml_w, ml_h,
         fill=CARD_BG, line_color=ACCENT2, line_width=Pt(2))
add_text(s3, "🤖 ML Service\n:5001 Flask",
         ml_x+Pt(5), ml_y+Pt(5), ml_w-Pt(10), Inches(0.45),
         size=12, bold=True, color=ACCENT2)
add_text(s3, "XGBoost ×4\nCatBoost ×1\n2–20ms inference\nIn-memory cache",
         ml_x+Pt(6), ml_y+Inches(0.52), ml_w-Pt(12), Inches(1.1),
         size=9, color=LIGHT_GRAY)
# Arrow Kafka → ML
for dy in [Inches(4.67), Inches(4.77), Inches(4.87)]:
    add_rect(s3, Inches(6.93), dy, Pt(3), Pt(3), fill=ACCENT2)

# ── Zookeeper ────────────────────────────────────────────────────
zoo_x = Inches(8.35); zoo_y = Inches(4.85)
add_rect(s3, zoo_x, zoo_y, Inches(2.25), Inches(0.7),
         fill=CARD_BG, line_color=ACCENT1, line_width=Pt(1.5))
add_text(s3, "Zookeeper :2181",
         zoo_x+Pt(5), zoo_y+Pt(5), Inches(2.1), Inches(0.55),
         size=11, bold=True, color=ACCENT1)
add_text(s3, "Kafka coordination",
         zoo_x+Pt(5), zoo_y+Inches(0.37), Inches(2.1), Inches(0.28),
         size=9, color=LIGHT_GRAY)

# ── 3 PostgreSQL Databases ───────────────────────────────────────
db_configs = [
    ("auth_db\n:5432", "users · onboarding\nverification_tokens", ACCENT1, Inches(0.25)),
    ("checkin_db\n:5433", "daily_check_ins · mood_logs\nsleep_logs · health_metrics\ndaily_ml_recommendation", ACCENT2, Inches(3.2)),
    ("noteai_db\n:5434", "journal_entries\nnotes · note_users", ACCENT4, Inches(7.1)),
]
db_y = Inches(4.75); db_w = Inches(2.7); db_h = Inches(1.6)
for (title, desc, clr, dx) in db_configs:
    add_rect(s3, dx, db_y, db_w, db_h,
             fill=rgb(0x0a,0x1f,0x30), line_color=clr, line_width=Pt(1.5))
    add_text(s3, f"🗄 {title}", dx+Pt(5), db_y+Pt(5), db_w-Pt(10), Inches(0.45),
             size=11, bold=True, color=clr)
    add_text(s3, "PostgreSQL 15", dx+Pt(5), db_y+Inches(0.48), db_w-Pt(10), Inches(0.25),
             size=8.5, color=ACCENT3, italic=True)
    add_text(s3, desc, dx+Pt(6), db_y+Inches(0.72), db_w-Pt(12), Inches(0.75),
             size=8.5, color=LIGHT_GRAY)

# ── Vertical arrows: Services → DBs ──────────────────────────────
arrow_pairs = [(Inches(4.5), Inches(1.75)), (Inches(7.0), Inches(4.55)), (Inches(9.5), Inches(6.35))]
# just dots
for ax, ay in [(Inches(1.6), Inches(3.2)), (Inches(1.6), Inches(3.36)), (Inches(1.6), Inches(3.52))]:
    add_rect(s3, ax, ay, Pt(3), Pt(3), fill=ACCENT1)

# Connection lines service → db (simple dotted lines)
def dot_line_v(slide, x, y1, y2, col, steps=5):
    for step in range(steps):
        dy = y1 + (y2-y1)*step/(steps-1)
        add_rect(slide, x, dy, Pt(3), Pt(3), fill=col)

dot_line_v(s3, Inches(4.5), Inches(3.22), Inches(4.75), ACCENT1)
dot_line_v(s3, Inches(7.0), Inches(3.22), Inches(4.75), ACCENT2)
dot_line_v(s3, Inches(9.5), Inches(3.22), Inches(4.75), ACCENT4)

# ── Right panel: docker info ─────────────────────────────────────
info_x = Inches(10.7); info_y = Inches(1.35)
add_rect(s3, info_x, info_y, Inches(2.5), Inches(5.95),
         fill=CARD_BG, line_color=ACCENT3, line_width=Pt(1.5))
add_text(s3, "🐳 Docker Compose", info_x+Pt(5), info_y+Pt(5), Inches(2.3), Inches(0.35),
         size=12, bold=True, color=ACCENT3)
add_text(s3, "10 контейнеров:", info_x+Pt(5), info_y+Inches(0.42), Inches(2.3), Inches(0.28),
         size=10, bold=True, color=WHITE)

docker_items = [
    ("NBAuthService", ACCENT1),
    ("NBCheckinService", ACCENT2),
    ("NoteAI-backend", ACCENT4),
    ("ML Flask Service", ACCENT2),
    ("PostgreSQL (×3)", ACCENT3),
    ("Apache Kafka", ACCENT3),
    ("Zookeeper", ACCENT1),
    ("pgAdmin UI", LIGHT_GRAY),
]
for di, (item, clr) in enumerate(docker_items):
    dy = info_y + Inches(0.75) + Inches(0.37)*di
    add_rect(s3, info_x+Pt(6), dy, Inches(2.2), Inches(0.3),
             fill=rgb(0x1a,0x2e,0x40), line_color=clr, line_width=Pt(1))
    add_text(s3, f"● {item}", info_x+Pt(10), dy+Pt(4), Inches(2.0), Inches(0.25),
             size=9.5, color=clr)

add_text(s3, "DigitalOcean Droplet\nProduction deploy",
         info_x+Pt(5), info_y+Inches(3.82), Inches(2.3), Inches(0.55),
         size=9.5, color=LIGHT_GRAY, italic=True)

add_text(s3, "HikariCP pool\nmax=10 conn/service",
         info_x+Pt(5), info_y+Inches(4.38), Inches(2.3), Inches(0.45),
         size=9.5, color=ACCENT1)

add_text(s3, "Flyway migrations\nV1__, V2__...",
         info_x+Pt(5), info_y+Inches(4.85), Inches(2.3), Inches(0.45),
         size=9.5, color=ACCENT3)

add_text(s3, "3 / 4", Inches(12.5), Inches(7.2), Inches(0.7), Inches(0.25),
         size=10, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════════
# SLIDE 4 — Backend Deep Dive: API · Security · Async · Gamification
# ════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank)
dark_bg(s4)
slide_header(s4,
    "Backend Deep Dive: API · Безопасность · Async · Геймификация",
    "Spring Boot 4 · Spring Security 6 · JWT HS256 · Kafka Events · XP System",
    accent=ACCENT4)

# ── Карточка 1: Security ─────────────────────────────────────────
c1x = Inches(0.25); c1y = Inches(1.3); cw4 = Inches(3.1); ch4 = Inches(2.9)
add_rect(s4, c1x, c1y, cw4, ch4, fill=CARD_BG, line_color=ACCENT4, line_width=Pt(2))
add_text(s4, "🔐 Безопасность",
         c1x+Pt(6), c1y+Pt(5), cw4-Pt(12), Inches(0.35),
         size=14, bold=True, color=ACCENT4)

sec_items = [
    ("JWT HS256",          "RFC 7519, accessToken 24ч",    ACCENT4),
    ("BCrypt cost=10",     "NIST SP 800-63B passwords",     ACCENT1),
    ("Refresh Token",      "UUID в БД, ротация при /refresh", ACCENT2),
    ("iOS Keychain",       "Secure Enclave, OWASP MASVS",   ACCENT3),
    ("Stateless auth",     "userId из request.getAttribute", ACCENT4),
    ("ML изолирован",      "Только внутри Docker-сети",     ACCENT1),
    ("BCrypt in AuthSvc",  "auth_db изолирован от остальных", ACCENT2),
]
for si, (title, desc, clr) in enumerate(sec_items):
    sy = c1y + Inches(0.42) + Inches(0.34)*si
    add_rect(s4, c1x+Pt(6), sy, cw4-Pt(12), Inches(0.28),
             fill=rgb(0x1a,0x2e,0x40))
    add_text(s4, f"▸ {title}:", c1x+Pt(10), sy+Pt(3), Inches(1.3), Inches(0.26),
             size=9.5, bold=True, color=clr)
    add_text(s4, desc, c1x+Pt(10)+Inches(1.32), sy+Pt(3), Inches(1.6), Inches(0.26),
             size=9, color=LIGHT_GRAY)

# ── Карточка 2: API Endpoints ────────────────────────────────────
c2x = Inches(3.5); c2y = Inches(1.3)
add_rect(s4, c2x, c2y, cw4, ch4, fill=CARD_BG, line_color=ACCENT1, line_width=Pt(2))
add_text(s4, "🌐 API Endpoints (40+)",
         c2x+Pt(6), c2y+Pt(5), cw4-Pt(12), Inches(0.35),
         size=14, bold=True, color=ACCENT1)

api_groups = [
    ("Auth (8081)",   "register · login · refresh\nonboarding · profile",    ACCENT1),
    ("Checkins",      "POST /checkins → Kafka trigger\n+ streak + XP auto",  ACCENT2),
    ("Sleep/Mood",    "POST /sleep → Kafka sleep.logged\nPOST /mood daily",  ACCENT3),
    ("Health ML",     "GET /health-metrics/today\nGET /ml/recommendations",   ACCENT4),
    ("Games",         "/brain-games, /game-sessions\n/new-game-sessions",     ACCENT1),
    ("NoteAI (8083)", "POST /journal → AI analyze\nPATCH auto-save debounce", ACCENT4),
    ("Swagger UI",    "/swagger-ui.html каждый сервис\n/apidocs (ML)",        LIGHT_GRAY),
]
for ai, (grp, desc, clr) in enumerate(api_groups):
    ay = c2y + Inches(0.42) + Inches(0.34)*ai
    add_rect(s4, c2x+Pt(6), ay, cw4-Pt(12), Inches(0.28), fill=rgb(0x1a,0x2e,0x40))
    add_text(s4, f"{grp}:", c2x+Pt(10), ay+Pt(3), Inches(1.15), Inches(0.26),
             size=9.5, bold=True, color=clr)
    add_text(s4, desc, c2x+Pt(10)+Inches(1.18), ay+Pt(3), Inches(1.75), Inches(0.26),
             size=8.5, color=LIGHT_GRAY)

# ── Карточка 3: Async Kafka ──────────────────────────────────────
c3x = Inches(6.75); c3y = Inches(1.3)
add_rect(s4, c3x, c3y, cw4, ch4, fill=CARD_BG, line_color=ACCENT3, line_width=Pt(2))
add_text(s4, "⚡ Kafka Async Flow",
         c3x+Pt(6), c3y+Pt(5), cw4-Pt(12), Inches(0.35),
         size=14, bold=True, color=ACCENT3)

flow_steps = [
    ("< 30 мс",  "POST /checkins → HTTP 201\n(только запись в БД)",  ACCENT3),
    ("async",    "Kafka: checkin.created\n→ HealthMetricsSaver",      ACCENT1),
    ("~200 мс",  "ML Service /predict\n→ HealthMetrics в БД",         ACCENT2),
    ("async",    "MLRecommendationConsumer\n→ /recommend/top3",        ACCENT4),
    ("fallback", "Если ML недоступен:\nформульный расчёт",             ACCENT4),
    ("topics",   "checkin.created · sleep.logged\ngame.completed",     ACCENT3),
    ("ordering", "key=userId → 1 partition\n→ гарантия порядка",       ACCENT1),
]
for fi, (timing, desc, clr) in enumerate(flow_steps):
    fy = c3y + Inches(0.42) + Inches(0.34)*fi
    add_rect(s4, c3x+Pt(6), fy, cw4-Pt(12), Inches(0.28), fill=rgb(0x1a,0x2e,0x40))
    badge_w = Inches(0.7)
    add_rect(s4, c3x+Pt(10), fy, badge_w, Inches(0.26), fill=clr)
    add_text(s4, timing, c3x+Pt(10), fy+Pt(2), badge_w, Inches(0.22),
             size=8, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_text(s4, desc, c3x+Pt(10)+badge_w+Pt(4), fy+Pt(3), Inches(2.0), Inches(0.26),
             size=8.5, color=LIGHT_GRAY)

# ── Карточка 4: Gamification + Tech Stack ────────────────────────
c4x = Inches(10.0); c4y = Inches(1.3)
add_rect(s4, c4x, c4y, cw4, ch4, fill=CARD_BG, line_color=ACCENT2, line_width=Pt(2))
add_text(s4, "🎮 Геймификация + Stack",
         c4x+Pt(6), c4y+Pt(5), cw4-Pt(12), Inches(0.35),
         size=14, bold=True, color=ACCENT2)

game_items = [
    ("XP formula",   "N²×100 per level\n(exp curve = retention)", ACCENT2),
    ("5 Tasks/day",  "AUTO: checkin·mood·sleep\ngame·journal → XP",  ACCENT3),
    ("Streaks",      "7d→×1.1 XP, 30d→×1.25 XP\nlast_date compare", ACCENT1),
    ("Adaptive",     "Когнит.скор→сложность игр\nEasy/Medium/Hard",   ACCENT4),
    ("Groq LLM",     "llama-3.3-70b temp=0.3\nSummarize·Analyze·Chat", ACCENT4),
    ("Java 17 LTS",  "G1 GC, sealed classes\nSpring AI integration",  ACCENT1),
    ("Monitoring",   "Swagger UI + pgAdmin\n:5050 + Docker logs",     LIGHT_GRAY),
]
for gi, (title, desc, clr) in enumerate(game_items):
    gy = c4y + Inches(0.42) + Inches(0.34)*gi
    add_rect(s4, c4x+Pt(6), gy, cw4-Pt(12), Inches(0.28), fill=rgb(0x1a,0x2e,0x40))
    add_text(s4, f"{title}:", c4x+Pt(10), gy+Pt(3), Inches(0.95), Inches(0.26),
             size=9.5, bold=True, color=clr)
    add_text(s4, desc, c4x+Pt(10)+Inches(0.97), gy+Pt(3), Inches(1.95), Inches(0.26),
             size=8.5, color=LIGHT_GRAY)

# ── Bottom row: metrics strip ────────────────────────────────────
strip_y = Inches(4.35)
add_rect(s4, 0, strip_y, W, Pt(2), fill=ACCENT1)

strip_items = [
    ("API latency", "< 30 мс", "(sync path)", ACCENT1),
    ("ML async", "~200 мс", "(Kafka consumer)", ACCENT3),
    ("XGBoost", "2–5 мс", "(inference)", ACCENT2),
    ("CatBoost", "10–20 мс", "(inference)", ACCENT4),
    ("BCrypt hash", "~100 мс", "(cost factor 10)", ACCENT4),
    ("JWT expiry", "86 400 с", "(24 часа)", ACCENT1),
    ("DB pool", "max 10", "(HikariCP)", ACCENT3),
    ("Containers", "10 total", "(Docker Compose)", ACCENT2),
]
strip_item_w = W / len(strip_items)
for si, (label, val, unit, clr) in enumerate(strip_items):
    sx = strip_item_w * si
    sy2 = strip_y + Pt(4)
    bg = CARD_BG if si%2==0 else rgb(0x1a,0x2e,0x40)
    add_rect(s4, sx, sy2, strip_item_w, Inches(1.0), fill=bg)
    add_text(s4, val, sx, sy2+Pt(6), strip_item_w, Inches(0.42),
             size=22, bold=True, color=clr, align=PP_ALIGN.CENTER)
    add_text(s4, unit, sx, sy2+Inches(0.42), strip_item_w, Inches(0.28),
             size=8.5, color=LIGHT_GRAY, align=PP_ALIGN.CENTER, italic=True)
    add_text(s4, label, sx, sy2+Inches(0.68), strip_item_w, Inches(0.28),
             size=9.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ── Bottom section: NoteAI + DB pattern ─────────────────────────
note_x = Inches(0.25); note_y = Inches(5.48)
add_rect(s4, note_x, note_y, Inches(6.2), Inches(1.78),
         fill=CARD_BG, line_color=ACCENT4, line_width=Pt(1.5))
add_text(s4, "🤖 NoteAI — Groq LLM Integration",
         note_x+Pt(6), note_y+Pt(5), Inches(6.0), Inches(0.32),
         size=12, bold=True, color=ACCENT4)
add_text(s4,
    "POST /journal → Spring AI (spring-ai-openai-starter v1.0.0-M6) → "
    "Groq Cloud API → llama-3.3-70b-versatile\n"
    "temperature=0.3 (детерминированные ответы)  ·  Groq LPU hardware → latency < 1 сек\n"
    "PATCH auto-save с debounce 2-3 сек  ·  Endpoints: /ai/summary · /ai/analyze · /ai/chat",
    note_x+Pt(7), note_y+Inches(0.4), Inches(6.0), Inches(1.3),
    size=10, color=LIGHT_GRAY)

db_x = Inches(6.6); db_y2 = Inches(5.48)
add_rect(s4, db_x, db_y2, Inches(6.5), Inches(1.78),
         fill=CARD_BG, line_color=ACCENT3, line_width=Pt(1.5))
add_text(s4, "🗄 Database-per-Service (Newman, 2015)",
         db_x+Pt(6), db_y2+Pt(5), Inches(6.3), Inches(0.32),
         size=12, bold=True, color=ACCENT3)
add_text(s4,
    "auth_db :5432  →  BCrypt passwords, JWT refresh tokens, onboarding\n"
    "checkin_db :5433  →  daily_check_ins · mood_logs · health_metrics · daily_ml_recommendation\n"
    "noteai_db :5434  →  journal_entries · notes\n"
    "Soft references по userId (no FK across DBs) · Flyway versioned migrations · HikariCP pool",
    db_x+Pt(7), db_y2+Inches(0.4), Inches(6.3), Inches(1.3),
    size=9.5, color=LIGHT_GRAY)

add_text(s4, "4 / 4", Inches(12.5), Inches(7.2), Inches(0.7), Inches(0.25),
         size=10, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════════
# Save
# ════════════════════════════════════════════════════════════════
out = '/Users/amangeldimadina/IdeaProjects/NeuroBalanceBackend/NeuroBalance_Presentation.pptx'
prs.save(out)
print(f'Saved: {out}')
